from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── helpers ──────────────────────────────────────────────────────────────────

def rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_textbox(slide, text, x, y, w, h, font_name, font_size, color,
                bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return txBox

def set_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_color)

def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width_pt=0, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex and line_width_pt > 0:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def stat_box(slide, x, y, w, h, fill_hex, border_hex, big_text, big_size, big_color,
             small_text, small_size='#475569'):
    # background rect
    bg = add_rect(slide, x, y, w, h, fill_hex, border_hex, 0)
    # left accent bar
    bar = add_rect(slide, x, y, 0.055, h, border_hex)
    # big number
    add_textbox(slide, big_text, x+0.12, y+0.08, w-0.2, 0.55,
                'Arial Black', big_size, big_color, bold=True)
    # small label
    add_textbox(slide, small_text, x+0.12, y+0.65, w-0.2, 0.75,
                'Arial', 10, '#475569', wrap=True)

def pill_box(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb('EFF6FF')
    shape.line.color.rgb = rgb('2563EB')
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'Arial'
    run.font.size = Pt(11)
    run.font.color.rgb = rgb('2563EB')

def step_box(slide, x, y, w, h, text):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb('F8FAFC')
    shape.line.color.rgb = rgb('2563EB')
    shape.line.width = Pt(0)
    # left bar
    bar = add_rect(slide, x, y, 0.045, h, '2563EB')
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = 'Arial'
    run.font.size = Pt(10.5)
    run.font.color.rgb = rgb('0F172A')

def add_table(slide, x, y, w, col_widths, rows_data, header_fill='0F172A'):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    cols = len(col_widths)
    nrows = len(rows_data)
    tbl = slide.shapes.add_table(nrows, cols, Inches(x), Inches(y),
                                  Inches(w), Inches(0.35 * nrows)).table
    # col widths
    total = sum(col_widths)
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)

    alt_fills = ['F8FAFC', 'FFFFFF']
    for r_idx, row_data in enumerate(rows_data):
        is_header = (r_idx == 0)
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell_text
            # clear auto-run
            for old_run in p.runs[:-1]:
                p._p.remove(old_run._r)
            run.font.name = 'Arial Black' if is_header else 'Arial'
            run.font.size = Pt(10.5 if is_header else 10)
            run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF) if is_header else RGBColor(0x0F,0x17,0x2A)
            # cell fill
            fill_hex = header_fill if is_header else alt_fills[r_idx % 2]
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr  = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', fill_hex)
    return tbl

# ── Presentation setup ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
blank_layout = prs.slide_layouts[6]  # blank

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, '0F172A')

# Shield shape (diamond as stand-in)
shield = s1.shapes.add_shape(4, Inches(4.4), Inches(0.35), Inches(1.2), Inches(1.3))
shield.fill.solid(); shield.fill.fore_color.rgb = rgb('2563EB')
shield.line.fill.background()

# Title
add_textbox(s1, 'AML Intelligence Platform',
            0.5, 1.9, 9, 0.7, 'Arial Black', 40, 'FFFFFF', bold=True, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(s1, 'Automated Graph-Based Money Laundering Detection',
            0.5, 2.7, 9, 0.45, 'Arial', 17, '93C5FD', align=PP_ALIGN.CENTER)

# Divider line
line_shape = s1.shapes.add_shape(1, Inches(2), Inches(3.3), Inches(6), Inches(0.02))
line_shape.fill.solid(); line_shape.fill.fore_color.rgb = rgb('2563EB')
line_shape.line.fill.background()

# Info lines
add_textbox(s1, 'iDEA 2.0  |  PS3: Fund Flow Tracking  |  Team Zeta',
            0.5, 3.48, 9, 0.3, 'Arial', 12, 'FFFFFF', align=PP_ALIGN.CENTER)
add_textbox(s1, 'Union Bank of India  x  K.J. Somaiya College of Engineering',
            0.5, 3.78, 9, 0.3, 'Arial', 11, '93C5FD', align=PP_ALIGN.CENTER)
add_textbox(s1, 'ideahackathon-1.onrender.com',
            0.5, 4.08, 9, 0.3, 'Arial', 11, '93C5FD', align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, 'FFFFFF')

add_textbox(s2, 'The Problem', 0.5, 0.25, 9, 0.65, 'Arial Black', 36, '0F172A', bold=True)
add_textbox(s2, 'Union Bank processes millions of transactions daily through NEFT, RTGS, and UPI. '
            'Fraud investigators have no automated system to trace how illicit funds move between accounts.',
            0.5, 0.95, 9, 0.55, 'Arial', 12, '475569', wrap=True)

# 3 stat boxes
stat_box(s2, 0.5,  1.65, 2.8, 1.5, 'EFF6FF', '2563EB', '3-5 Days', 30, '2563EB',
         'to manually reconstruct one transaction trail')
stat_box(s2, 3.45, 1.65, 2.8, 1.5, 'FEF2F2', 'DC2626', '7 Days', 30, 'DC2626',
         'FIU deadline for Suspicious Transaction Reports')
stat_box(s2, 6.4,  1.65, 2.8, 1.5, 'FFFBEB', 'D97706', 'Rs.2000-5000 Cr', 22, 'D97706',
         'lost annually to laundering (FATF 2023)')

add_textbox(s2, 'Rule-based threshold systems are bypassed by structuring. '
            'Manual investigation is reactive - funds exit before detection. '
            'Money laundering is a graph problem. No existing system looks at the graph.',
            0.5, 3.38, 9.0, 0.9, 'Arial', 11, '475569', wrap=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — HOW IT WORKS
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, 'FFFFFF')

add_textbox(s3, 'How It Works', 0.5, 0.25, 9, 0.65, 'Arial Black', 36, '0F172A', bold=True)

# Left column
add_textbox(s3, 'Detection Pipeline', 0.5, 1.05, 5.3, 0.3, 'Arial Black', 12, '2563EB', bold=True)

steps = [
    '1.  IBM HI-Small Dataset (500K transactions) loaded into NetworkX directed graph',
    '2.  Labelled Pipeline classifies 8 AML pattern types: Cycle, Fan-Out, Fan-In, Scatter-Gather, Gather-Scatter, Bipartite, Stack, Random',
    '3.  Unlabelled Pipeline scores all accounts on 7 behavioural signals - accounts with 2+ signals are flagged',
    '4.  Random Forest ML Model scores each alert on 16 graph features. Precision: 83%',
]
sy = 1.42
for step in steps:
    step_box(s3, 0.5, sy, 5.3, 0.73, step)
    sy += 0.83

# Right column
add_textbox(s3, 'Tech Stack', 6.2, 1.05, 3.4, 0.3, 'Arial Black', 12, '2563EB', bold=True)

pills = ['Python + FastAPI', 'NetworkX Graph Engine', 'scikit-learn Random Forest',
         'Cytoscape.js Visualisation', 'Chart.js Dashboard', 'Deployed on Render', 'IBM HI-Small Dataset']
py = 1.42
for pill in pills:
    pill_box(s3, 6.2, py, 3.3, 0.38, pill)
    py += 0.46

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RESULTS
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, 'FFFFFF')

add_textbox(s4, 'What We Detected', 0.5, 0.25, 9, 0.65, 'Arial Black', 36, '0F172A', bold=True)

# 2x2 stat grid
def result_box(slide, x, y, big, label, border):
    bg = add_rect(slide, x, y, 2.2, 1.45, 'FFFFFF', border, 0)
    bar = add_rect(slide, x, y, 0.055, 1.45, border)
    add_textbox(slide, big, x+0.15, y+0.1, 1.9, 0.7,
                'Arial Black', 38, border, bold=True)
    add_textbox(slide, label, x+0.15, y+0.82, 1.9, 0.55,
                'Arial', 9, '475569', wrap=True)

result_box(s4, 0.50, 1.1, '676',   'Alerts Detected', '2563EB')
result_box(s4, 2.85, 1.1, '3,091', 'Unlabelled Suspicious Accounts', '059669')
result_box(s4, 0.50, 2.7, '8',     'AML Pattern Types', 'D97706')
result_box(s4, 2.85, 2.7, '83%',   'ML Precision Score', 'DC2626')

# Comparison table
rows = [
    ['BUILT (POC)', 'PLANNED (PRODUCTION)'],
    ['Graph detection pipeline', 'Real-time Kafka stream ingestion'],
    ['8 AML pattern classifiers', '15+ pattern library'],
    ['ML fraud scoring', 'Graph Neural Network model'],
    ['Investigation dashboard', 'RBI-compliant STR report export'],
    ['Whitelist suppression', 'KYC/KYB entity integration'],
]
add_table(s4, 5.1, 1.1, 4.65, [2.1, 2.55], rows)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TEAM
# ════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, 'FFFFFF')

add_textbox(s5, 'Team Zeta', 0.5, 0.25, 9, 0.65, 'Arial Black', 36, '0F172A', bold=True)

team_rows = [
    ['MEMBER', 'CONTRIBUTION'],
    ['Viraj Sanghavi', 'Backend, graph pipeline, ML model, frontend, deployment'],
    ['Sonal', 'Testing and validation'],
    ['Archit', 'Dataset research and analysis'],
    ['Suruchi', 'Documentation and presentation'],
]
add_table(s5, 0.5, 1.15, 9.0, [2.5, 6.5], team_rows)

# Built vs Planned
add_textbox(s5, 'Built:', 0.5, 3.42, 0.7, 0.3, 'Arial Black', 11, '2563EB', bold=True)
add_textbox(s5, 'Full detection pipeline, ML model, dashboard, validation, deployment',
            1.25, 3.42, 8.2, 0.3, 'Arial', 11, '0F172A')

add_textbox(s5, 'Planned:', 0.5, 3.82, 0.9, 0.3, 'Arial Black', 11, '475569', bold=True)
add_textbox(s5, 'Kafka streaming, Graph Neural Networks, RBI STR reports, role-based access control',
            1.45, 3.82, 8.0, 0.3, 'Arial', 11, '475569')

add_textbox(s5, 'K.J. Somaiya College of Engineering  |  iDEA 2.0 Phase 2  |  PS3: Fund Flow Tracking',
            0.5, 4.55, 9, 0.35, 'Arial', 10, '475569', align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DELIVERABLES
# ════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, '0F172A')

add_textbox(s6, 'All Deliverables', 0.5, 0.2, 9, 0.6, 'Arial Black', 32, 'FFFFFF',
            bold=True, align=PP_ALIGN.CENTER)

boxes = [
    (0.4,  1.1,  'D1 - Problem + Solution Brief', '[Add Google Doc link]'),
    (5.1,  1.1,  'D2 - Technical Demo Video',      '[Add YouTube Unlisted link]'),
    (0.4,  2.35, 'D3 - Architecture Document',     '[Add Google Doc link]'),
    (5.1,  2.35, 'D4 - GitHub Repository',         'github.com/VirajSanghavi007/IdeaHackathon'),
    (0.4,  3.6,  'D5 - Pitch Video',               '[Add YouTube Unlisted link]'),
    (5.1,  3.6,  'Live Application',               'ideahackathon-1.onrender.com'),
]

for bx, by, label, link in boxes:
    bg = add_rect(s6, bx, by, 4.5, 1.05, '1E293B', '2563EB', 0.8)
    add_textbox(s6, label, bx+0.15, by+0.08, 4.2, 0.38,
                'Arial Black', 10.5, '93C5FD', bold=True)
    add_textbox(s6, link,  bx+0.15, by+0.55, 4.2, 0.38,
                'Arial', 9.5, 'FFFFFF')

add_textbox(s6, 'iDEA 2.0 Phase 2  |  PS3: Fund Flow Tracking  |  Team Zeta  |  Union Bank of India',
            0.5, 5.05, 9, 0.3, 'Arial', 9, '475569', align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r'C:\Users\viraj\Downloads\Projects\IdeaHackathon\D5_Pitch_Deck.pptx'
prs.save(out)
print(f'Saved: {out}')
